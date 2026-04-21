"""
Generate docs/agentic-workflow.pptx from the agentic workflow diagram.

Requires:
  python3 -m pip install python-pptx

Usage: python3 scripts/gen-workflow-ppt.py
"""

try:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches, Pt
except ModuleNotFoundError as exc:
    if exc.name == "pptx":
        raise SystemExit(
            "Missing dependency 'python-pptx'. Install it with: python3 -m pip install python-pptx"
        ) from exc
    raise

from pathlib import Path

# ── Colour palette ──────────────────────────────────────────────────────────
BG          = RGBColor(0x0d, 0x11, 0x17)
CARD_BG     = RGBColor(0x16, 0x1b, 0x22)
BORDER_BLUE = RGBColor(0x4a, 0x90, 0xd9)
TEXT_BLUE   = RGBColor(0x93, 0xc5, 0xfd)
SUB_BLUE    = RGBColor(0x5a, 0x8a, 0xbd)
BORDER_PRP  = RGBColor(0xc9, 0x7d, 0xd4)
FILL_PRP    = RGBColor(0x3d, 0x1a, 0x3d)
TEXT_PRP    = RGBColor(0xe0, 0xa8, 0xe8)
TITLE_PRP   = RGBColor(0xfa, 0xe8, 0xfd)
FILL_GRN    = RGBColor(0x12, 0x2b, 0x1e)
BORDER_GRN  = RGBColor(0x4c, 0xaf, 0x7d)
TEXT_GRN    = RGBColor(0xd1, 0xfa, 0xe5)
SUB_GRN     = RGBColor(0x86, 0xef, 0xac)
ORC_FILL    = RGBColor(0x2a, 0x1e, 0x00)
ORC_BORDER  = RGBColor(0xd4, 0xa0, 0x17)
ORC_TEXT    = RGBColor(0xd4, 0xa0, 0x17)
WHITE       = RGBColor(0xe6, 0xed, 0xf3)
ARROW       = RGBColor(0x4a, 0x90, 0xd9)
CONNECTOR   = RGBColor(0x21, 0x26, 0x2d)

# ── Slide dimensions: 16:9 widescreen ────────────────────────────────────────
W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

SCRIPT_DIR = Path(__file__).resolve().parent
WORKSPACE_ROOT = SCRIPT_DIR.parent
DOCS_DIR = WORKSPACE_ROOT / "docs"

# ── Helpers ──────────────────────────────────────────────────────────────────

def add_rect(slide, x, y, w, h, fill_rgb, border_rgb, border_pt=2.0):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_rgb
    line = shape.line
    line.color.rgb = border_rgb
    line.width = Pt(border_pt)
    return shape

def set_text(shape, lines, align=PP_ALIGN.CENTER):
    """lines = list of (text, font_size, bold, rgb)"""
    tf = shape.text_frame
    tf.word_wrap = False
    tf.auto_size = None
    from pptx.util import Pt as _Pt
    for i, (text, fsize, bold, rgb) in enumerate(lines):
        if i == 0:
            para = tf.paragraphs[0]
        else:
            para = tf.add_paragraph()
        para.alignment = align
        run = para.add_run()
        run.text = text
        run.font.size = _Pt(fsize)
        run.font.bold = bold
        run.font.color.rgb = rgb
    # Vertical centre
    tf.margin_top    = Pt(0)
    tf.margin_bottom = Pt(0)
    tf.margin_left   = Pt(4)
    tf.margin_right  = Pt(4)

def centre_v(shape):
    from pptx.enum.text import MSO_ANCHOR
    shape.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

def add_label(slide, x, y, w, h, text, fsize, bold, rgb):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = False
    para = tf.paragraphs[0]
    para.alignment = PP_ALIGN.CENTER
    run = para.add_run()
    run.text = text
    run.font.size = Pt(fsize)
    run.font.bold = bold
    run.font.color.rgb = rgb
    return tb

def add_line(slide, x1, y1, x2, y2, rgb, width_pt=2.0):
    """Draw a solid line as a thin filled rectangle (Google Slides compatible)."""
    thick = Pt(width_pt)
    if abs(x2 - x1) >= abs(y2 - y1):  # horizontal
        lx = min(x1, x2)
        ly = min(y1, y2) - thick // 2
        lw = abs(x2 - x1)
        lh = max(thick, Pt(1))
    else:  # vertical
        lx = min(x1, x2) - thick // 2
        ly = min(y1, y2)
        lw = max(thick, Pt(1))
        lh = abs(y2 - y1)
    s = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, lx, ly, lw, lh)
    s.fill.solid()
    s.fill.fore_color.rgb = rgb
    s.line.fill.background()  # no border on the rect itself
    return s

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ═══════════════════════════════════════════════════════════════════════════════
slide1_layout = prs.slide_layouts[6]  # blank
slide1 = prs.slides.add_slide(slide1_layout)

# Background
bg = add_rect(slide1, 0, 0, W, H, BG, BG, 0)

# Accent bar (left edge)
add_rect(slide1, 0, 0, Inches(0.12), H, BORDER_BLUE, BORDER_BLUE, 0)

# Title text
tb = slide1.shapes.add_textbox(Inches(0.8), Inches(2.8), Inches(11), Inches(1.2))
tf = tb.text_frame
para = tf.paragraphs[0]
para.alignment = PP_ALIGN.LEFT
run = para.add_run()
run.text = "Agentic Workflow"
run.font.size = Pt(52)
run.font.bold = True
run.font.color.rgb = WHITE

sub_tb = slide1.shapes.add_textbox(Inches(0.8), Inches(3.9), Inches(11), Inches(0.6))
sub_tf = sub_tb.text_frame
sub_para = sub_tf.paragraphs[0]
sub_para.alignment = PP_ALIGN.LEFT
sub_run = sub_para.add_run()
sub_run.text = "Orchestrated gate-based AI development pipeline"
sub_run.font.size = Pt(20)
sub_run.font.color.rgb = RGBColor(0x8b, 0x94, 0x9e)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Full Workflow Diagram
# ═══════════════════════════════════════════════════════════════════════════════
slide2_layout = prs.slide_layouts[6]
slide2 = prs.slides.add_slide(slide2_layout)
add_rect(slide2, 0, 0, W, H, BG, BG, 0)

# Slide title
add_label(slide2, Inches(0.4), Inches(0.18), Inches(12), Inches(0.5),
          "Agentic Workflow", 24, True, WHITE)

# ── Layout constants ────────────────────────────────────────────────────────
# 9 blocks across rows, with a U-turn connector on the right
# Slide usable area: x 0.3→13.0  y 0.8→7.2
# Block size
BW = Inches(1.55)   # block width
BH = Inches(0.9)    # block height

# Row 1 y-top
R1Y = Inches(1.05)
# Row 2 y-top
R2Y = Inches(5.2)
# Orchestrator bar
OY  = Inches(3.3)
OH  = Inches(0.38)

# Row 1 x positions (5 blocks, left→right)
R1X = [Inches(0.28), Inches(2.15), Inches(4.02), Inches(5.89), Inches(7.76)]
# Row 2 x positions (6 blocks, right→left visually but stored left→right for coords)
# right-most is Build, left-most is Done
R2X = [Inches(0.28), Inches(2.15), Inches(4.02), Inches(5.89), Inches(7.76), Inches(9.63)]
# Row 2 logical order right→left: Build(9.63), CodeReview(7.76), RuntimeQA(5.89), MergeReady(4.02), POMerge(2.15), Done(0.28)

# ── Row 1 blocks ─────────────────────────────────────────────────────────────
blocks_r1 = [
    # (x, fill, border, lines)
    (R1X[0], FILL_PRP, BORDER_PRP, [("👤 Product", 12, True, TITLE_PRP), ("Owner", 10, False, TEXT_PRP)]),
    (R1X[1], CARD_BG,  BORDER_BLUE, [("Requirement", 11, False, TEXT_BLUE), ("Challenger", 9, False, SUB_BLUE)]),
    (R1X[2], CARD_BG,  BORDER_BLUE, [("PRD", 11, False, TEXT_BLUE), ("PRD Agent", 9, False, SUB_BLUE)]),
    (R1X[3], CARD_BG,  BORDER_BLUE, [("UX+Design", 11, False, TEXT_BLUE), ("Orchestrator", 9, False, SUB_BLUE)]),
    (R1X[4], CARD_BG,  BORDER_BLUE, [("Architecture", 11, False, TEXT_BLUE), ("Arch Agent", 9, False, SUB_BLUE)]),
]

for (bx, fill, border, lines) in blocks_r1:
    s = add_rect(slide2, bx, R1Y, BW, BH, fill, border)
    set_text(s, lines)
    centre_v(s)

# ── Row 2 blocks (right→left: Build, CodeReview, RuntimeQA, MergeReady, POMerge, Done) ──
blocks_r2 = [
    (R2X[5], CARD_BG,  BORDER_BLUE, [("Build", 11, False, TEXT_BLUE), ("Dev Agent", 9, False, SUB_BLUE)]),
    (R2X[4], CARD_BG,  BORDER_BLUE, [("Code Review", 11, False, TEXT_BLUE), ("Code Reviewer", 9, False, SUB_BLUE)]),
    (R2X[3], CARD_BG,  BORDER_BLUE, [("Runtime QA", 11, False, TEXT_BLUE), ("Runtime QA Agent", 9, False, SUB_BLUE)]),
    (R2X[2], CARD_BG,  BORDER_BLUE, [("Merge Ready", 11, False, TEXT_BLUE), ("Gate 6", 9, False, SUB_BLUE)]),
    (R2X[1], FILL_PRP, BORDER_PRP,  [("👤 PO", 12, True, TITLE_PRP), ("Merges PR", 10, False, TEXT_PRP), ("PO-only", 9, False, RGBColor(0xb0, 0x7c, 0xc0))]),
    (R2X[0], FILL_GRN, BORDER_GRN,  [("✅ Done", 13, True, TEXT_GRN), ("Shipped", 10, False, SUB_GRN)]),
]

for (bx, fill, border, lines) in blocks_r2:
    s = add_rect(slide2, bx, R2Y, BW, BH, fill, border)
    set_text(s, lines)
    centre_v(s)

# ── Row 1 horizontal arrows ───────────────────────────────────────────────────
GAP = Inches(0.08)
for i in range(4):
    x1 = R1X[i] + BW + GAP
    x2 = R1X[i+1] - GAP
    y  = R1Y + BH / 2
    add_line(slide2, x1, y, x2, y, ARROW, 2.0)

# ── U-turn connector: right edge of Architecture → right margin → down → Build right edge ──
# Architecture right edge
arch_right = R1X[4] + BW
build_right = R2X[5] + BW
MARGIN_X = Inches(11.6)
R1_MID = R1Y + BH / 2
R2_MID = R2Y + BH / 2

# Segment 1: horizontal right from Architecture centre-right to margin
add_line(slide2, arch_right + Inches(0.05), R1_MID, MARGIN_X, R1_MID, ARROW, 2.0)
# Segment 2: vertical down margin
add_line(slide2, MARGIN_X, R1_MID, MARGIN_X, R2_MID, ARROW, 2.0)
# Segment 3: horizontal left from margin to Build right edge
add_line(slide2, MARGIN_X, R2_MID, build_right + Inches(0.05), R2_MID, ARROW, 2.0)

# ── Row 2 horizontal arrows (right→left) ─────────────────────────────────────
r2_order = [R2X[5], R2X[4], R2X[3], R2X[2], R2X[1], R2X[0]]
for i in range(5):
    x1 = r2_order[i] - GAP
    x2 = r2_order[i+1] + BW + GAP
    y  = R2Y + BH / 2
    add_line(slide2, x1, y, x2, y, ARROW, 2.0)

# ── Orchestrator swim-lane bar ────────────────────────────────────────────────
OX1 = R1X[1]            # starts at Requirement
OX2 = R2X[5] + BW       # ends at Build right edge
orc_bar = add_rect(slide2, OX1, OY, OX2 - OX1, OH, ORC_FILL, ORC_BORDER, 1.5)
set_text(orc_bar, [("🎯  Orchestrator — supervises every gate", 10, True, ORC_TEXT)])
centre_v(orc_bar)

# Connector ticks: row-1 block bottoms → bar top
ORC_TOP    = OY
ORC_BOTTOM = OY + OH

for bx in R1X[1:]:    # Gates 1-4
    cx = bx + BW / 2
    add_line(slide2, cx, R1Y + BH, cx, ORC_TOP,    ORC_BORDER, 1.2)

# Connector ticks: bar bottom → row-2 block tops
for bx in [R2X[5], R2X[4], R2X[3], R2X[2]]:   # Build, CodeReview, RuntimeQA, MergeReady
    cx = bx + BW / 2
    add_line(slide2, cx, ORC_BOTTOM, cx, R2Y,      ORC_BORDER, 1.2)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Gate legend
# ═══════════════════════════════════════════════════════════════════════════════
slide3_layout = prs.slide_layouts[6]
slide3 = prs.slides.add_slide(slide3_layout)
add_rect(slide3, 0, 0, W, H, BG, BG, 0)
add_label(slide3, Inches(0.4), Inches(0.18), Inches(12), Inches(0.5),
          "Gate Sequence at a Glance", 24, True, WHITE)

gates = [
    ("Gate 1", "Requirement", "Requirement Challenger grills the brief; exposes ambiguity before work starts."),
    ("Gate 2", "PRD",         "PRD Agent converts validated requirements into a structured product spec."),
    ("Gate 3", "Design",      "Orchestrator executes UX+Design in Gate 3A; Design QA Agent validates fidelity in Gate 3B."),
    ("Gate 4", "Architecture","Architecture Agent maps modules, boundaries, risks, and task breakdown."),
    ("Gate 5", "Build",       "Dev Agent implements and tests against acceptance criteria on a feature branch."),
    ("Code Review", "",       "Code Reviewer (Copilot) reviews the PR; feedback loop until 0 comments."),
    ("Gate 5.5", "Runtime QA","Runtime QA Agent validates live behaviour across viewport × theme matrix."),
    ("Gate 6", "Merge Ready", "Orchestrator confirms all evidence; Product Owner merges the PR."),
]

ROW_H = Inches(0.62)
START_Y = Inches(0.9)

for i, (gate, title, desc) in enumerate(gates):
    y = START_Y + i * ROW_H
    # Gate badge
    badge = add_rect(slide3, Inches(0.3), y + Inches(0.06), Inches(1.55), Inches(0.42), CARD_BG, BORDER_BLUE, 1.5)
    set_text(badge, [(gate, 9, True, TEXT_BLUE)])
    centre_v(badge)
    # Title
    if title:
        add_label(slide3, Inches(2.05), y + Inches(0.08), Inches(1.8), Inches(0.4),
                  title, 11, True, WHITE)
    # Description
    dtb = slide3.shapes.add_textbox(Inches(3.9), y + Inches(0.06), Inches(9.0), Inches(0.48))
    dtf = dtb.text_frame
    dtf.word_wrap = True
    dpara = dtf.paragraphs[0]
    drun = dpara.add_run()
    drun.text = desc
    drun.font.size = Pt(11)
    drun.font.color.rgb = RGBColor(0x8b, 0x94, 0x9e)

# ── Save ──────────────────────────────────────────────────────────────────────
DOCS_DIR.mkdir(parents=True, exist_ok=True)
out = DOCS_DIR / "agentic-workflow.pptx"
prs.save(out)
print(f"Saved → {out.relative_to(WORKSPACE_ROOT)}")
